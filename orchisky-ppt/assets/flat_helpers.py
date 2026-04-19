# -*- coding: utf-8 -*-
"""
扁平化渲染工具库（原生 PPT 专用）
消除 PowerPoint 默认的 UI 按鈕感和立体感。
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

try:
    from constants import FONT_SANS, FONT_SERIF
except ImportError:
    FONT_SANS = 'Microsoft YaHei'
    FONT_SERIF = 'Songti SC'


def _kill_theme_inheritance(shape):
    """'移除 shape 的主题样式引用，确保完全扁平"""
    sp = shape._element
    for style in sp.findall(qn('p:style')):
        sp.remove(style)
    spPr = sp.find('.//' + qn('p:spPr'))
    if spPr is None:
        spPr = sp.find('.//' + qn('a:spPr'))
    if spPr is not None:
        for eff in spPr.findall(qn('a:effectLst')):
            spPr.remove(eff)
        for eff in spPr.findall(qn('a:effectRef')):
            spPr.remove(eff)
    return shape


def add_flat_rect(slide, x, y, w, h, fill=None, rx=0):
    """扁平矩形，默认直角、无边框、无主题"""
    shape_type = MSO_SHAPE.RECTANGLE if rx == 0 else MSO_SHAPE.ROUNDED_RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(x / 96.0), Inches(y / 96.0),
        Inches(w / 96.0), Inches(h / 96.0)
    )
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.fill.background()
    if rx > 0 and shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        shape.adjustments[0] = rx / min(w, h) if min(w, h) > 0 else 0.1
    _kill_theme_inheritance(shape)
    return shape


def add_flat_line(slide, x1, y1, x2, y2,
                  color=RGBColor(0x1A, 0x1A, 0x1A),
                  width_pt=1.0, dash=None):
    """扁平线条"""
    line = slide.shapes.add_connector(
        1,
        Inches(x1 / 96.0), Inches(y1 / 96.0),
        Inches(x2 / 96.0), Inches(y2 / 96.0)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    if dash is not None:
        line.line.dash_style = dash
    _kill_theme_inheritance(line)
    return line


def add_flat_circle(slide, cx, cy, r, fill=None, line_color=None, line_width_pt=0):
    """扁平圆"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches((cx - r) / 96.0), Inches((cy - r) / 96.0),
        Inches((2 * r) / 96.0), Inches((2 * r) / 96.0)
    )
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color is not None and line_width_pt > 0:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    _kill_theme_inheritance(shape)
    return shape


def _set_font_faces(rPr, font_name):
    for tag in ('a:ea', 'a:latin', 'a:cs'):
        existing = rPr.find(qn(tag))
        if existing is not None:
            rPr.remove(existing)
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', font_name)
    latin = etree.SubElement(rPr, qn('a:latin'))
    latin.set('typeface', font_name)
    cs = etree.SubElement(rPr, qn('a:cs'))
    cs.set('typeface', font_name)


def add_flat_text(slide, x, y, w, h, text,
                  font_size=12, bold=False, italic=False,
                  color=RGBColor(0x1A, 0x1A, 0x1A),
                  font=FONT_SANS,
                  align='left', anchor='top',
                  letter_spacing=0):
    """扁平文本框"""
    tb = slide.shapes.add_textbox(
        Inches(x / 96.0), Inches(y / 96.0),
        Inches(w / 96.0), Inches(h / 96.0)
    )
    tb.fill.background()
    tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    anchor_map = {'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM, 'top': MSO_ANCHOR.TOP}
    tf.vertical_anchor = anchor_map.get(anchor, MSO_ANCHOR.TOP)
    align_map = {'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT, 'left': PP_ALIGN.LEFT}
    p = tf.paragraphs[0]
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    if letter_spacing > 0:
        rPr.set('spc', str(letter_spacing * 100))
    _set_font_faces(rPr, font)
    _kill_theme_inheritance(tb)
    return tb


def add_page_header(slide, breadcrumb):
    add_flat_text(slide, 40, 32, 800, 20, breadcrumb,
                  font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA), letter_spacing=2)


def add_page_footer(slide, source, page_no):
    add_flat_text(slide, 40, 688, 800, 20, source,
                  font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA))
    add_flat_text(slide, 1100, 688, 140, 20, str(page_no),
                  font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA), align='right')


def add_action_title_simple(slide, title, y=82):
    add_flat_text(slide, 40, y, 1200, 30, title,
                  font_size=22, bold=True,
                  color=RGBColor(0x1A, 0x1A, 0x1A), font=FONT_SERIF)


def add_action_title_standard(slide, title, y=36):
    add_flat_rect(slide, 40, y, 3, 52, fill=RGBColor(0x00, 0x3D, 0x7A))
    add_flat_text(slide, 52, y + 14, 1200, 30, title,
                  font_size=22, bold=True,
                  color=RGBColor(0x1A, 0x1A, 0x1A), font=FONT_SERIF)
    add_flat_line(slide, 40, y + 58, 1240, y + 58,
                  color=RGBColor(0xE2, 0xE2, 0xE2), width_pt=0.5)


def add_top_brand_bar(slide):
    add_flat_rect(slide, 0, 0, 1280, 3, fill=RGBColor(0x00, 0x3D, 0x7A))
