"""Core PPTX assembly: create_pptx_with_native_svg."""

from __future__ import annotations

import hashlib
import re
import shutil
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from .drawingml_converter import convert_svg_to_slide_shapes
from .pptx_dimensions import get_format, FORMATS


def _append_relationship(rels_path: Path, rel_type: str, target: str) -> str:
    with open(rels_path, 'r', encoding='utf-8') as f:
        rels_content = f.read()
    rid_numbers = [int(match) for match in re.findall(r'Id="rId(\d+)"', rels_content)]
    next_rid = f'rId{max(rid_numbers, default=0) + 1}'
    rel_xml = f'  <Relationship Id="{next_rid}" Type="{rel_type}" Target="{target}"/>'
    rels_content = rels_content.replace('</Relationships>', rel_xml + '\n</Relationships>')
    with open(rels_path, 'w', encoding='utf-8') as f:
        f.write(rels_content)
    return next_rid


def create_pptx_with_native_svg(
    svg_files: list,
    output_path: Path,
    canvas_format: str = 'ppt169',
    verbose: bool = True,
    transition=None,
    transition_duration: float = 0.5,
    auto_advance=None,
    use_compat_mode: bool = False,
    notes: dict = None,
    enable_notes: bool = True,
    use_native_shapes: bool = True,
) -> bool:
    """Create a PPTX file by converting SVG files to native DrawingML shapes.

    Args:
        svg_files: List of SVG Path objects.
        output_path: Output PPTX path.
        canvas_format: Slide format key ('ppt169', 'ppt43', etc.).
        verbose: Whether to print conversion progress.
        transition: Transition effect name (None to disable).
        use_compat_mode: Ignored when use_native_shapes=True.
        notes: Dict mapping SVG stem to Markdown notes text.
        enable_notes: Whether to embed speaker notes.
        use_native_shapes: Convert SVG to native DrawingML (recommended).

    Returns:
        True if all slides converted successfully.
    """
    if not svg_files:
        print("Error: No SVG files found")
        return False

    fmt = get_format(canvas_format)
    width_emu = Emu(fmt.width)
    height_emu = Emu(fmt.height)

    if verbose:
        print(f"  Slide format: {fmt.name} ({fmt.width} x {fmt.height} EMU)")
        print(f"  SVG file count: {len(svg_files)}")
        print(f"  Mode: {'Native DrawingML shapes' if use_native_shapes else 'SVG embedding'}")
        print()

    temp_dir = Path(tempfile.mkdtemp())

    try:
        # Create base PPTX
        prs = Presentation()
        prs.slide_width = width_emu
        prs.slide_height = height_emu
        blank_layout = prs.slide_layouts[6]
        for _ in svg_files:
            prs.slides.add_slide(blank_layout)
        base_pptx = temp_dir / 'base.pptx'
        prs.save(str(base_pptx))

        # Extract PPTX
        extract_dir = temp_dir / 'pptx_content'
        with zipfile.ZipFile(base_pptx, 'r') as zf:
            zf.extractall(extract_dir)

        media_dir = extract_dir / 'ppt' / 'media'
        media_dir.mkdir(exist_ok=True)

        success_count = 0
        has_any_image = False
        media_cache: dict = {}
        notes_slides_created: set = set()

        for i, svg_path in enumerate(svg_files, 1):
            slide_num = i
            try:
                if use_native_shapes:
                    slide_xml, media_files_dict, rel_entries = convert_svg_to_slide_shapes(
                        svg_path, slide_num=slide_num, verbose=verbose,
                    )

                    # Write slide XML
                    slide_xml_path = extract_dir / 'ppt' / 'slides' / f'slide{slide_num}.xml'
                    with open(slide_xml_path, 'w', encoding='utf-8') as f:
                        f.write(slide_xml)

                    # Write media files
                    media_name_map: dict = {}
                    for media_name, media_data in media_files_dict.items():
                        ext = media_name.rsplit('.', 1)[-1].lower()
                        media_hash = hashlib.sha256(media_data).hexdigest()
                        cache_key = (ext, media_hash)
                        cached_name = media_cache.get(cache_key)
                        if cached_name is None:
                            cached_name = media_name
                            media_cache[cache_key] = cached_name
                            with open(media_dir / cached_name, 'wb') as f:
                                f.write(media_data)
                        media_name_map[media_name] = cached_name

                    for rel in rel_entries:
                        target = rel.get('target', '')
                        if target.startswith('../media/'):
                            media_name = target.split('../media/', 1)[1]
                            mapped = media_name_map.get(media_name)
                            if mapped:
                                rel['target'] = f'../media/{mapped}'

                    # Build rels
                    rels_dir = extract_dir / 'ppt' / 'slides' / '_rels'
                    rels_dir.mkdir(exist_ok=True)
                    rels_path = rels_dir / f'slide{slide_num}.xml.rels'

                    extra_rels = ''
                    for rel in rel_entries:
                        extra_rels += (
                            f'\n  <Relationship Id="{rel["id"]}" '
                            f'Type="{rel["type"]}" Target="{rel["target"]}"/>'
                        )

                    rels_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>{extra_rels}
</Relationships>'''
                    with open(rels_path, 'w', encoding='utf-8') as f:
                        f.write(rels_xml)

                    for media_name in media_files_dict:
                        ext = media_name.rsplit('.', 1)[-1].lower()
                        if ext in ('png', 'jpg', 'jpeg', 'gif', 'webp'):
                            has_any_image = True

                # Notes handling
                if enable_notes and notes:
                    svg_stem = svg_path.stem
                    notes_text = notes.get(svg_stem, '')
                    if notes_text:
                        try:
                            from .pptx_notes import md_to_plain, build_notes_slide_xml
                            plain_text = md_to_plain(notes_text)
                        except Exception:
                            plain_text = notes_text

                        notes_slides_dir = extract_dir / 'ppt' / 'notesSlides'
                        notes_slides_dir.mkdir(exist_ok=True)
                        notes_xml_path = notes_slides_dir / f'notesSlide{slide_num}.xml'
                        try:
                            from .pptx_notes import build_notes_slide_xml
                            notes_xml_bytes = build_notes_slide_xml(plain_text)
                            with open(notes_xml_path, 'wb') as f:
                                f.write(notes_xml_bytes)
                        except Exception:
                            pass
                        notes_slides_created.add(slide_num)

                if verbose:
                    mode_str = " (Native)" if use_native_shapes else " (SVG)"
                    notes_str = " +notes" if slide_num in notes_slides_created else ""
                    print(f"  [{i}/{len(svg_files)}] {svg_path.name}{mode_str}{notes_str}")

                success_count += 1

            except Exception as e:
                if verbose:
                    print(f"  [{i}/{len(svg_files)}] {svg_path.name} - Error: {e}")

        # Update [Content_Types].xml
        content_types_path = extract_dir / '[Content_Types].xml'
        with open(content_types_path, 'r', encoding='utf-8') as f:
            content_types = f.read()

        types_to_add = []
        if has_any_image:
            for ext, ct in [('png', 'image/png'), ('jpg', 'image/jpeg'), ('jpeg', 'image/jpeg')]:
                if f'Extension="{ext}"' not in content_types:
                    types_to_add.append(f'  <Default Extension="{ext}" ContentType="{ct}"/>')

        if notes_slides_created:
            for idx in sorted(notes_slides_created):
                override = (
                    f'  <Override PartName="/ppt/notesSlides/notesSlide{idx}.xml" '
                    f'ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>'
                )
                if override not in content_types:
                    types_to_add.append(override)

        if types_to_add:
            content_types = content_types.replace(
                '</Types>', '\n'.join(types_to_add) + '\n</Types>',
            )
            with open(content_types_path, 'w', encoding='utf-8') as f:
                f.write(content_types)

        # Repackage PPTX
        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for file_path in extract_dir.rglob('*'):
                if file_path.is_file():
                    arcname = file_path.relative_to(extract_dir)
                    zf.write(file_path, arcname)

        if verbose:
            print()
            print(f"[Done] Saved: {output_path}")
            print(f"  Succeeded: {success_count}, Failed: {len(svg_files) - success_count}")

        return success_count == len(svg_files)

    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)
