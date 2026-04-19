# -*- coding: utf-8 -*-
"""
⚠️ LEGACY · V3.0 老路径资产 · V4.0 不主动加载
如需使用，请通过 svg_to_pptx_wrapper.py 调用。
"""

from pptx import Presentation
from pptx.util import Pt
import copy

def fix_theme_fonts(input_path: str, output_path: str = None) -> None:
    """修复 PPTX 中的主题字体映射，确保中文字体正确显示。"""
    if output_path is None:
        output_path = input_path
    
    prs = Presentation(input_path)
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name in ('Songti SC', 'PingFang SC', None):
                            run.font.name = 'Microsoft YaHei'
    
    prs.save(output_path)
    print(f"Fixed fonts: {output_path}")

if __name__ == '__main__':
    import sys
    if len(sys.argv) >= 2:
        fix_theme_fonts(sys.argv[1], sys.argv[2] if len(sys.argv) >= 3 else None)
